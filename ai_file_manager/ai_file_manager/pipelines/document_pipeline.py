"""
pipelines/document_pipeline.py — Extract text from document files
(.txt, .md, .csv, .docx, .pptx, .xlsx, .html, .json, .xml, etc.)
"""

from __future__ import annotations
import asyncio
import json
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Optional

from core.router import ProcessingResult
from utils.config import load_config
from utils.logger import get_logger

log = get_logger("document_pipeline")


async def process_document(file_info: dict) -> ProcessingResult:
    """Pipeline entry point for generic document files."""
    path = Path(file_info["path"])
    ext = path.suffix.lower()

    loop = asyncio.get_event_loop()

    try:
        text = await loop.run_in_executor(None, _extract_document, path, ext)
    except Exception as exc:
        log.error(f"Document pipeline error for {path.name}: {exc}")
        return ProcessingResult(file_info, error=str(exc))

    if not text:
        return ProcessingResult(file_info, error="No text extracted from document")

    max_chars = int(load_config().processing.document.max_chars)
    truncated = text[:max_chars]
    full_text = f"[Document: {path.name}]\n\n{truncated}"

    return ProcessingResult(file_info, text=full_text)


def _extract_document(path: Path, ext: str) -> str:
    """Dispatch to the appropriate text extractor for the file extension."""
    extractors = {
        ".txt":  _read_plain_text,
        ".md":   _read_plain_text,
        ".rst":  _read_plain_text,
        ".csv":  _read_csv,
        ".tsv":  _read_csv,
        ".json": _read_json,
        ".yaml": _read_plain_text,
        ".yml":  _read_plain_text,
        ".xml":  _read_xml,
        ".html": _read_html,
        ".htm":  _read_html,
        ".docx": _read_docx,
        ".doc":  _read_docx,
        ".pptx": _read_pptx,
        ".xlsx": _read_xlsx,
        ".xls":  _read_xlsx,
        ".rtf":  _read_plain_text,
    }
    fn = extractors.get(ext, _read_plain_text)
    return fn(path)


def _read_plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _read_csv(path: Path) -> str:
    """Read first N rows of CSV as a text summary."""
    try:
        import csv
        rows = []
        with open(path, encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 100:
                    rows.append(f"... ({i} rows total)")
                    break
                rows.append(", ".join(str(c) for c in row))
        return "\n".join(rows)
    except Exception:
        return _read_plain_text(path)


def _read_json(path: Path) -> str:
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)[:8000]
    except Exception:
        return _read_plain_text(path)


def _read_xml(path: Path) -> str:
    try:
        tree = ET.parse(str(path))
        root = tree.getroot()
        parts = []

        def _walk(node, depth=0):
            text = (node.text or "").strip()
            if text:
                parts.append("  " * depth + f"<{node.tag}>: {text}")
            for child in node:
                _walk(child, depth + 1)

        _walk(root)
        return "\n".join(parts)[:6000]
    except Exception:
        return _read_plain_text(path)


def _read_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
        html = path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        # Remove script/style
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)[:6000]
    except ImportError:
        # Fallback: crude tag removal via regex
        import re
        html = path.read_text(encoding="utf-8", errors="replace")
        return re.sub(r"<[^>]+>", " ", html)[:6000]
    except Exception:
        return _read_plain_text(path)


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except ImportError:
        log.warning("python-docx not installed: pip install python-docx")
        return f"[DOCX: {path.name} — install python-docx for content extraction]"
    except Exception as exc:
        log.warning(f"DOCX read failed for {path.name}: {exc}")
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"[Slide {i}]: " + " | ".join(slide_texts))
        return "\n".join(parts)
    except ImportError:
        log.warning("python-pptx not installed: pip install python-pptx")
        return f"[PPTX: {path.name} — install python-pptx for content extraction]"
    except Exception as exc:
        log.warning(f"PPTX read failed for {path.name}: {exc}")
        return ""


def _read_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= 50:
                    rows.append("...")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)[:6000]
    except ImportError:
        log.warning("openpyxl not installed: pip install openpyxl")
        return f"[XLSX: {path.name} — install openpyxl for content extraction]"
    except Exception as exc:
        log.warning(f"XLSX read failed for {path.name}: {exc}")
        return ""
